from pathlib import Path

from pptx import Presentation

from src.markdown_parser import Deck, Slide
from src.ppt_builder import build_presentation
from src.theme import get_theme


def test_semantic_design_layouts_render(tmp_path: Path):
    deck = Deck(
        title="施工方案汇报",
        subtitle="设计引擎测试",
        slides=[
            Slide(title="汇报重点", bullets=["高边坡开挖支护属于重点风险控制内容"], layout="highlight"),
            Slide(title="汇报目录", bullets=["工程概况", "施工部署", "风险控制"], layout="agenda"),
            Slide(title="工程概况", body=["洞口高边坡开挖支护，包含截水沟、防护网、锚喷支护。"], layout="overview"),
            Slide(title="施工工艺流程", bullets=["测量放样", "开挖", "支护", "验收"], layout="process"),
            Slide(title="风险辨识与控制", bullets=["边坡坍塌", "物体打击", "触电伤害"], layout="risk"),
            Slide(title="安全保障措施", bullets=["交底到位", "专人检查", "闭环整改"], layout="checklist"),
        ],
    )
    output_path = tmp_path / "semantic.pptx"

    build_presentation(deck, output_path, theme=get_theme("construction"))

    prs = Presentation(output_path)
    assert len(prs.slides) == 7
    assert output_path.stat().st_size > 0


def test_semantic_layouts_keep_text_density_under_control(tmp_path: Path):
    long_item = "这是一个非常长的施工安全控制要求，用于验证页面设计引擎不会把过多文字硬塞进一个很小的文本框导致字体溢出。"
    deck = Deck(
        title="容量控制测试",
        subtitle="设计引擎测试",
        slides=[
            Slide(title="风险辨识与控制", bullets=[long_item] * 6, layout="risk"),
            Slide(title="安全保障措施", bullets=[long_item] * 7, layout="checklist"),
            Slide(title="施工工艺流程", bullets=[long_item] * 6, layout="process"),
        ],
    )
    output_path = tmp_path / "density.pptx"

    build_presentation(deck, output_path, theme=get_theme("construction"))

    prs = Presentation(output_path)
    for index in range(1, len(prs.slides)):
        slide = prs.slides[index]
        text_length = 0
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text_length += sum(len(paragraph.text) for paragraph in shape.text_frame.paragraphs)
        assert text_length < 360


def test_process_layout_avoids_placeholders_and_duplicate_step_text(tmp_path: Path):
    deck = Deck(
        title="流程页测试",
        subtitle="",
        slides=[Slide(title="施工流程", bullets=["测量放样", "分层开挖"], layout="process")],
    )
    output_path = tmp_path / "process.pptx"

    build_presentation(deck, output_path, theme=get_theme("construction"))

    prs = Presentation(output_path)
    slide = prs.slides[1]
    text = "\n".join(
        paragraph.text
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
        for paragraph in shape.text_frame.paragraphs
    )
    assert "待补充" not in text
    assert text.count("测量放样") == 1


def test_builder_does_not_emit_ellipsis_for_shortened_text(tmp_path: Path):
    deck = Deck(
        title="省略号测试",
        subtitle="",
        slides=[
            Slide(
                title="工程概况",
                body=["湖北秭归抽水蓄能电站位于湖北省宜昌市秭归县茅坪镇，上水库沿村村通公路、芝茅公路至秭归县。"],
                layout="summary",
            )
        ],
    )
    output_path = tmp_path / "ellipsis.pptx"

    build_presentation(deck, output_path, theme=get_theme("construction"))

    prs = Presentation(output_path)
    text = "\n".join(
        paragraph.text
        for shape in prs.slides[1].shapes
        if getattr(shape, "has_text_frame", False)
        for paragraph in shape.text_frame.paragraphs
    )
    assert "..." not in text
    assert "…" not in text
